import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import os
import sys
import numpy as np
from scipy.spatial import ConvexHull
from scipy.spatial import QhullError  # Správný import

from pptx import Presentation
from pptx.util import Inches

print("--- START SKRIPTU ---")

# --- Konfigurace ---
DATA_FILE_PATH = "makroseis2025.xlsx"
SHEET_NAME = 0
EQ_YEAR_TARGET = 2025
EQ_MONTH_TARGET = 4
EQ_DAY_TARGET = 24
EQ_DATETIME_UTC_STR = (
    f"{EQ_YEAR_TARGET}-{EQ_MONTH_TARGET:02d}-{EQ_DAY_TARGET:02d} 17:32:47.3"
)
EQ_LAT = 49.422
EQ_LON = 14.043
EQ_MAGNITUDE = 3.1
EQ_GFU_ID = 1773
EQ_LOCATION_NAME = "u Mirotic"
TIME_WINDOW_HOURS_FILTER = 1.5
CENTER_LAT_CR_ZOOMED = 49.81746
CENTER_LON_CR_ZOOMED = 15.47490
ZOOM_LEVEL_CR_ZOOMED = 6.0
OUTPUT_DIR = f"analyza_vysledky_{EQ_LOCATION_NAME.lower().replace(' ', '_').replace('-', '_')}_{EQ_YEAR_TARGET}"
print(f"Výstupní adresář bude: {os.path.abspath(OUTPUT_DIR)}")

COL_OBS_DATETIME = "eqdatetime"
COL_LAT = "lat"
COL_LON = "lon"
COL_IN_BUILDING = "pozorovaniodkud"
COL_FEAR = "reakcepanika"
COL_TREMOR_TYPE = "popispohybu"
COL_FELT_BY = "kolikpozorvenku"
COLS_OBJECT_MOVEMENT_DETAILS = [
    "nabytektezky",
    "okna",
    "dvere",
    "zavespredmety",
    "nadobi",
    "malepredmety",
    "kapalina",
]
COL_DAMAGE_OVERALL = "poskozomitka"
COLS_SOUNDS = [
    "zvukzavrzani",
    "zvukduneni",
    "zvukvibrace",
    "zvukhuceni",
    "zvukvrzani",
    "zvukvitr",
    "zvukexploze",
    "zvuktanku",
    "zvuknevin",
]
NEGATIVE_OR_EMPTY_VALUES = ["", "0", "ne", "no", "false", "nan", "null"]

# --- Načtení a základní příprava dat ---
print(f"\n--- Načítání dat z: {DATA_FILE_PATH} ---")
try:
    df = pd.read_excel(
        DATA_FILE_PATH, sheet_name=SHEET_NAME, na_values=["NULL", "null", ""]
    )
    print(f"Úspěšně načteno {len(df)} řádků z Excelu.")
    if df.empty:
        print("CHYBA: Načtený DataFrame je prázdný.")
        sys.exit("Skript ukončen - prázdný DataFrame.")
except FileNotFoundError:
    print(f"CHYBA: Soubor {DATA_FILE_PATH} nebyl nalezen.")
    sys.exit(f"Skript ukončen - soubor nenalezen.")
except Exception as e:
    print(f"CHYBA při načítání Excelu: {e}")
    sys.exit(f"Skript ukončen - chyba Excelu: {e}")

print(f"\n--- Zpracování časových údajů (sloupec '{COL_OBS_DATETIME}') ---")
try:
    if COL_OBS_DATETIME not in df.columns:
        print(f"CHYBA: Sloupec '{COL_OBS_DATETIME}' nenalezen.")
        sys.exit(f"Skript ukončen - chybí {COL_OBS_DATETIME}.")
    df[COL_OBS_DATETIME] = pd.to_datetime(
        df[COL_OBS_DATETIME], dayfirst=True, errors="coerce"
    )
    df.dropna(subset=[COL_OBS_DATETIME], inplace=True)
    print(f"Po konverzi času: {len(df)} řádků.")
    if df.empty:
        print("CHYBA: Žádná platná časová data.")
        sys.exit("Skript ukončen - žádná časová data.")
    if df[COL_OBS_DATETIME].dt.tz is None:
        df[COL_OBS_DATETIME] = (
            df[COL_OBS_DATETIME]
            .dt.tz_localize("Europe/Prague", ambiguous="infer")
            .dt.tz_convert("UTC")
        )
    else:
        df[COL_OBS_DATETIME] = df[COL_OBS_DATETIME].dt.tz_convert("UTC")
    print("Časová zóna aplikována.")
except Exception as e:
    print(f"CHYBA při zpracování času: {e}")
    sys.exit(f"Skript ukončen - chyba času: {e}")

print(f"\n--- Zpracování souřadnic ('{COL_LAT}', '{COL_LON}') ---")
try:
    if COL_LAT not in df.columns or COL_LON not in df.columns:
        print(f"CHYBA: Sloupce souřadnic nenalezeny.")
        sys.exit("Skript ukončen - chybí souřadnice.")
    df[COL_LAT] = pd.to_numeric(df[COL_LAT], errors="coerce")
    df[COL_LON] = pd.to_numeric(df[COL_LON], errors="coerce")
    df.dropna(subset=[COL_LAT, COL_LON], inplace=True)
    print(f"Po konverzi souřadnic: {len(df)} řádků.")
    if df.empty:
        print("CHYBA: Žádná platná data souřadnic.")
        sys.exit("Skript ukončen - žádné souřadnice.")
except Exception as e:
    print(f"CHYBA při zpracování souřadnic: {e}")
    sys.exit(f"Skript ukončen - chyba souřadnic: {e}")

print(f"\n--- Filtrování ({EQ_YEAR_TARGET}/{EQ_MONTH_TARGET}) ---")
df_filtered_year_month = df[
    (df[COL_OBS_DATETIME].dt.year == EQ_YEAR_TARGET)
    & (df[COL_OBS_DATETIME].dt.month == EQ_MONTH_TARGET)
].copy()
print(
    f"Nalezeno {len(df_filtered_year_month)} záznamů pro {EQ_MONTH_TARGET}/{EQ_YEAR_TARGET}."
)
if df_filtered_year_month.empty:
    sys.exit(f"Skript ukončen - žádné záznamy pro {EQ_MONTH_TARGET}/{EQ_YEAR_TARGET}.")

target_eq_datetime_utc = pd.Timestamp(EQ_DATETIME_UTC_STR, tz="UTC")
time_delta_filter = pd.Timedelta(hours=TIME_WINDOW_HOURS_FILTER)
start_time_filter = target_eq_datetime_utc - time_delta_filter
end_time_filter = target_eq_datetime_utc + time_delta_filter
print(f"Cílový čas (UTC): {target_eq_datetime_utc}")
print(f"Časové okno: {start_time_filter} do {end_time_filter} (UTC)")
df_event = df_filtered_year_month.loc[
    (df_filtered_year_month[COL_OBS_DATETIME] >= start_time_filter)
    & (df_filtered_year_month[COL_OBS_DATETIME] <= end_time_filter)
].copy()
print(
    f"\n--- ANALÝZA PRO {EQ_LOCATION_NAME} ({target_eq_datetime_utc.strftime('%Y-%m-%d %H:%M:%S %Z')}) ---"
)
print(f"Nalezeno {len(df_event)} pozorování v okně +/- {TIME_WINDOW_HOURS_FILTER}h.")
if df_event.empty:
    sys.exit("Skript ukončen - žádné záznamy v okně události.")

print(f"\n--- Kontrola výstupního adresáře: {OUTPUT_DIR} ---")
if not os.path.exists(OUTPUT_DIR):
    try:
        os.makedirs(OUTPUT_DIR)
        print(f"Vytvořen adresář: {OUTPUT_DIR}")
    except Exception as e:
        sys.exit(f"Skript ukončen - chyba vytváření adresáře: {e}")
else:
    print(f"Adresář {OUTPUT_DIR} již existuje.")


# Pomocná funkce pro hovertemplate
def build_hovertemplate_string(hover_data_config, main_hover_col):
    template_parts = []
    # Najdeme index hlavního sloupce pro customdata
    custom_data_keys = list(hover_data_config.keys())

    if main_hover_col and main_hover_col in custom_data_keys:
        main_col_idx = custom_data_keys.index(main_hover_col)
        template_parts.append(
            f"<b>{main_hover_col}</b>: %{{customdata[{main_col_idx}]}}<br>"
        )

    for i, col_name in enumerate(custom_data_keys):
        if hover_data_config[col_name] and col_name != main_hover_col:
            template_parts.append(f"{col_name}: %{{customdata[{i}]}}<br>")
    template_parts.append("<extra></extra>")
    return "".join(template_parts)


# --- Helper funkce pro tvorbu map ---
def create_custom_map(
    df_map_data,
    color_column_name,
    map_title_suffix,
    output_filename_base,
    category_orders_dict=None,
    color_discrete_map_dict=None,
    hover_data_extra=None,
    hover_text_column_name=None,
    show_isoseismal_areas=False,
    ems_color_map_for_hulls=None,
    sort_ems_key_func=None,
):
    if color_column_name not in df_map_data.columns and color_column_name is not None:
        print(
            f"INFO: Sloupec '{color_column_name}' pro mapu '{map_title_suffix}' nenalezen."
        )
        return None
    png_path = os.path.join(
        OUTPUT_DIR, f"mapa_{output_filename_base}_{EQ_YEAR_TARGET}.png"
    )

    # Příprava hover dat
    active_hover_data_config = {
        COL_OBS_DATETIME: True
    }  # COL_OBS_DATETIME je vždy přítomen
    # Efektivní sloupec pro hlavní text v hoveru (tučně)
    effective_main_hover_text_col = (
        hover_text_column_name
        if hover_text_column_name and hover_text_column_name in df_map_data.columns
        else color_column_name
    )

    if (
        effective_main_hover_text_col
        and effective_main_hover_text_col in df_map_data.columns
    ):
        active_hover_data_config[effective_main_hover_text_col] = True
    if hover_data_extra:
        for item in hover_data_extra:
            if item in df_map_data.columns:
                active_hover_data_config[item] = True

    # Odstraníme duplicity a zajistíme, že klíče v active_hover_data_config jsou unikátní a existují v df_map_data
    final_hover_data_cols = [
        col
        for col in list(dict.fromkeys(active_hover_data_config.keys()))
        if col in df_map_data.columns
    ]
    df_for_customdata = df_map_data[final_hover_data_cols]

    if (
        show_isoseismal_areas
        and color_column_name == "EMS_Intensity_Est"
        and "EMS_Intensity_Est" in df_map_data.columns
        and ems_color_map_for_hulls
        and sort_ems_key_func
    ):
        fig = go.Figure()
        valid_ems_levels_for_hulls_asc = sorted(
            [
                level
                for level in ems_color_map_for_hulls.keys()
                if "Neklasifikováno" not in level and "Nepocítěno" not in level
            ],
            key=sort_ems_key_func,
            reverse=False,
        )
        df_for_hulls = df_map_data.copy()
        df_for_hulls[COL_LAT] = pd.to_numeric(df_for_hulls[COL_LAT], errors="coerce")
        df_for_hulls[COL_LON] = pd.to_numeric(df_for_hulls[COL_LON], errors="coerce")
        df_for_hulls.dropna(
            subset=[COL_LAT, COL_LON, "EMS_Intensity_Est"], inplace=True
        )
        if "EMS_Intensity_Est_SortVal" not in df_for_hulls.columns:
            ems_string_to_sort_val = {
                level: sort_ems_key_func(level)
                for level in df_for_hulls["EMS_Intensity_Est"].unique()
            }
            df_for_hulls = df_for_hulls.assign(
                EMS_Intensity_Est_SortVal=df_for_hulls["EMS_Intensity_Est"].map(
                    ems_string_to_sort_val
                )
            )
        for ems_level_str in valid_ems_levels_for_hulls_asc:
            current_level_sort_val = sort_ems_key_func(ems_level_str)
            points_for_hull_df = df_for_hulls[
                df_for_hulls["EMS_Intensity_Est_SortVal"] >= current_level_sort_val
            ]
            if len(points_for_hull_df) >= 3:
                coordinates = points_for_hull_df[[COL_LON, COL_LAT]].values
                try:
                    hull = ConvexHull(coordinates)
                    hull_lons = coordinates[hull.vertices, 0]
                    hull_lats = coordinates[hull.vertices, 1]
                    hull_lons = np.append(hull_lons, hull_lons[0])
                    hull_lats = np.append(hull_lats, hull_lats[0])
                    fig.add_trace(
                        go.Scattermapbox(
                            lon=hull_lons,
                            lat=hull_lats,
                            mode="none",
                            fill="toself",
                            fillcolor=ems_color_map_for_hulls.get(
                                ems_level_str, "grey"
                            ),
                            name=f"Oblast {ems_level_str.split(' - ')[0]}",
                            hoverinfo="name",
                            legendgroup="isoseismals",
                            showlegend=True,
                            opacity=0.5,
                        )
                    )
                except QhullError:
                    print(
                        f"INFO: ConvexHull pro {ems_level_str} ({len(points_for_hull_df)} bodů) přeskočen."
                    )
                except Exception as e:
                    print(f"CHYBA: ConvexHull pro {ems_level_str}: {e}")

        point_colors_mapped = None
        if (
            color_column_name
            and color_discrete_map_dict
            and color_column_name in df_map_data.columns
        ):
            point_colors_mapped = (
                df_map_data[color_column_name]
                .map(color_discrete_map_dict)
                .fillna("grey")
            )  # grey pro nemapované

        fig.add_trace(
            go.Scattermapbox(
                lat=df_map_data[COL_LAT],
                lon=df_map_data[COL_LON],
                mode="markers",
                marker=go.scattermapbox.Marker(
                    size=8,
                    color=point_colors_mapped
                    if point_colors_mapped is not None
                    else "blue",  # Použije namapované barvy
                    opacity=0.9,
                ),
                customdata=df_for_customdata.values,
                hovertemplate=build_hovertemplate_string(
                    dict.fromkeys(final_hover_data_cols, True),
                    effective_main_hover_text_col,
                ),
                name="Pozorování",
                legendgroup="observations",
                showlegend=True,
            )
        )
    else:
        fig = px.scatter_mapbox(
            df_map_data,
            lat=COL_LAT,
            lon=COL_LON,
            hover_name=df_map_data.index
            if df_map_data.index.name
            else None,  # Může být jméno sloupce nebo index
            hover_data={
                col: True for col in final_hover_data_cols
            },  # Použijeme dict pro hover_data
            color=color_column_name
            if color_column_name and color_column_name in df_map_data.columns
            else None,
            category_orders=category_orders_dict
            if color_column_name and color_column_name in df_map_data.columns
            else None,
            color_discrete_map=color_discrete_map_dict
            if color_column_name and color_column_name in df_map_data.columns
            else None,
            size_max=10,
            opacity=0.8,
            zoom=ZOOM_LEVEL_CR_ZOOMED,
            center={"lat": CENTER_LAT_CR_ZOOMED, "lon": CENTER_LON_CR_ZOOMED},
        )
        if not (color_column_name and color_column_name in df_map_data.columns):
            fig.update_traces(marker=dict(size=7))

    fig.add_trace(
        go.Scattermapbox(
            lat=[EQ_LAT],
            lon=[EQ_LON],
            mode="markers",
            marker=go.scattermapbox.Marker(
                size=17, color="red", symbol="star", opacity=1
            ),
            name=f"Epicentrum (Mag: {EQ_MAGNITUDE})",
            text=[
                f"Epicentrum {EQ_LOCATION_NAME}<br>Magnituda: {EQ_MAGNITUDE}<br>ID: {EQ_GFU_ID}"
            ],
            hoverinfo="text",
            showlegend=True,
        )
    )
    main_title = f"Zemětřesení {EQ_LOCATION_NAME}"
    legend_title_for_map = map_title_suffix
    if (
        not (color_column_name and color_column_name in df_map_data.columns)
        and not show_isoseismal_areas
    ):
        legend_title_for_map = "Legenda"
    fig.update_layout(
        mapbox_style="open-street-map",
        margin={"r": 10, "t": 50, "l": 10, "b": 10},
        title=main_title,
        legend_title_text=legend_title_for_map,
        legend=dict(
            bgcolor="rgba(255,255,255,0.9)",
            bordercolor="Black",
            borderwidth=1,
            title_font_family="Arial",
            font=dict(family="Arial", size=10, color="black"),
            itemsizing="constant",
            traceorder="reversed",
            yanchor="top",
            y=0.98,
            xanchor="right",
            x=0.98,
        ),
        mapbox_zoom=ZOOM_LEVEL_CR_ZOOMED,
        mapbox_center={"lat": CENTER_LAT_CR_ZOOMED, "lon": CENTER_LON_CR_ZOOMED},
    )
    try:
        html_path = os.path.join(
            OUTPUT_DIR, f"mapa_{output_filename_base}_{EQ_YEAR_TARGET}.html"
        )
        fig.write_html(html_path)
        print(f"Mapa '{map_title_suffix}' uložena do HTML: {html_path}")
        fig.write_image(png_path, scale=3, width=1000, height=750)
        print(f"Mapa '{map_title_suffix}' uložena do PNG: {png_path}")
        return png_path
    except Exception as e:
        print(f"CHYBA při ukládání mapy '{map_title_suffix}': {e}.")
        if "kaleido" in str(e).lower():
            print("      Nainstalujte 'kaleido': pip install kaleido")
        return None


# ZDE JE DEFINICE create_bar_chart PŘESUNUTA NA SPRÁVNÉ MÍSTO
# --- Helper funkce pro tvorbu sloupcových grafů ---
def create_bar_chart(
    data_series,
    chart_title,
    filename_base,
    xaxis_title,
    yaxis_title="Počet pozorování",
    show_percentages=False,
):
    if data_series.empty:
        print(f"INFO: Graf '{chart_title}' se negeneruje (žádná data).")
        return None  # Vracíme None, pokud se graf negeneruje

    total_observations_for_chart_title = data_series.sum()
    fig = px.bar(
        x=data_series.index,
        y=data_series.values,
        labels={"x": xaxis_title, "y": yaxis_title},
        title=chart_title
        + f" (celkem {total_observations_for_chart_title} pozorování)",
    )
    fig.update_layout(xaxis_title=xaxis_title, yaxis_title=yaxis_title)

    if show_percentages and not df_event.empty:
        percentages = (data_series / len(df_event)) * 100
        fig.update_traces(
            texttemplate="%{y} (%{customdata:.1f}%)",
            textposition="outside",
            customdata=percentages,
        )
    else:
        fig.update_traces(texttemplate="%{y}", textposition="outside")
    try:
        path = os.path.join(OUTPUT_DIR, f"graf_{filename_base}_{EQ_YEAR_TARGET}.html")
        fig.write_html(path)
        print(f"Graf '{chart_title}' uložen do: {path}")
        # Mohli bychom vracet i cestu k HTML grafu, pokud by to bylo potřeba
        # png_path_graf = os.path.join(OUTPUT_DIR, f"graf_{filename_base}_{EQ_YEAR_TARGET}.png")
        # fig.write_image(png_path_graf)
        # print(f"Graf '{chart_title}' uložen do PNG: {png_path_graf}")
        return fig  # Vracíme objekt figury
    except Exception as e:
        print(f"CHYBA při ukládání grafu '{chart_title}': {e}")
        if "kaleido" in str(e).lower():
            print("      Nainstalujte 'kaleido': pip install kaleido")
        return None


# --- Předběžné zpracování kategorií ---
print("\n--- Předzpracování kategorií ---")
df_event.loc[:, "Mist_Pozorovani_Kat_Full"] = (
    df_event[COL_IN_BUILDING]
    .fillna("")
    .astype(str)
    .str.lower()
    .apply(lambda x: "Doma (uvnitř)" if x == "budova" else "Venku/Nezadáno v budově")
)
df_event.loc[:, "Mist_Pozorovani_Legenda"] = (
    df_event[COL_IN_BUILDING]
    .fillna("")
    .astype(str)
    .str.lower()
    .apply(lambda x: "Doma" if x == "budova" else "Venku/Nezadáno")
)
df_event.loc[:, "Pocit_Kategorie_Text_Full"] = (
    df_event[COL_FELT_BY]
    .astype(str)
    .str.lower()
    .str.strip()
    .map(
        {
            "pouze vy": "Pocítil(a) jen respondent",
            "většina ano": "Pocítila většina přítomných",
        }
    )
    .fillna("Nezadáno/Jiná odpověď")
)
df_event.loc[:, "Pocit_Kategorie_Legenda"] = (
    df_event[COL_FELT_BY]
    .astype(str)
    .str.lower()
    .str.strip()
    .map({"pouze vy": "Jen respondent", "většina ano": "Většina"})
    .fillna("Nezadáno")
)
df_event.loc[:, "Intenzita_Kat_Full"] = (
    df_event[COL_TREMOR_TYPE].astype(str).str.lower().fillna("Nezadáno")
)
df_event.loc[:, "Intenzita_Kat_Legenda"] = df_event["Intenzita_Kat_Full"]
df_event.loc[:, "Strach_Pocit_Kat_Text_Full"] = (
    (pd.to_numeric(df_event[COL_FEAR], errors="coerce") == 1)
    .map({True: "Ano (strach/panika)", False: "Ne/Nezadáno strach"})
    .fillna("Ne/Nezadáno strach")
)
df_event.loc[:, "Strach_Pocit_Legenda"] = (
    (pd.to_numeric(df_event[COL_FEAR], errors="coerce") == 1)
    .map({True: "Ano", False: "Ne/Nezadáno"})
    .fillna("Ne/Nezadáno")
)
actual_movement_detail_cols = [
    col for col in COLS_OBJECT_MOVEMENT_DETAILS if col in df_event.columns
]
if actual_movement_detail_cols:
    df_event.loc[:, "Pohyb_Predmetu_Agregovany_Bool"] = df_event[
        actual_movement_detail_cols
    ].apply(
        lambda r: any(
            pd.notna(v) and str(v).strip().lower() not in NEGATIVE_OR_EMPTY_VALUES
            for v in r
        ),
        axis=1,
    )
    df_event.loc[:, "Pohyb_Predmetu_Agregovany_Text_Full"] = df_event[
        "Pohyb_Predmetu_Agregovany_Bool"
    ].map({True: "Ano (pohyb předmětů)", False: "Ne (žádný pohyb předmětů)"})
    df_event.loc[:, "Pohyb_Predmetu_Legenda"] = df_event[
        "Pohyb_Predmetu_Agregovany_Bool"
    ].map({True: "Ano", False: "Ne"})
else:
    df_event.loc[:, "Pohyb_Predmetu_Agregovany_Text_Full"] = (
        "Nezadáno (info o pohybu chybí)"
    )
    df_event.loc[:, "Pohyb_Predmetu_Legenda"] = "Nezadáno"
df_event.loc[:, "Poskozeni_Obecne_Text_Full"] = (
    df_event[COL_DAMAGE_OVERALL]
    .astype(str)
    .str.lower()
    .str.strip()
    .map({"bylo": "Ano (poškození hlášeno)", "nebylo": "Ne (poškození nehlášeno)"})
    .fillna("Nezadáno/Jiná hodnota")
)
df_event.loc[:, "Poskozeni_Obecne_Legenda"] = (
    df_event[COL_DAMAGE_OVERALL]
    .astype(str)
    .str.lower()
    .str.strip()
    .map({"bylo": "Ano", "nebylo": "Ne"})
    .fillna("Nezadáno")
)
actual_sound_cols = [col for col in COLS_SOUNDS if col in df_event.columns]
if actual_sound_cols:
    df_event.loc[:, "Zvuk_Reportovan_Bool"] = df_event[actual_sound_cols].apply(
        lambda r: any(
            pd.notna(v) and str(v).strip().lower() not in NEGATIVE_OR_EMPTY_VALUES
            for v in r
        ),
        axis=1,
    )
    df_event.loc[:, "Zvuk_Reportovan_Text_Full"] = df_event["Zvuk_Reportovan_Bool"].map(
        {True: "Ano (zvuk reportován)", False: "Ne (bez zvuku)"}
    )
    df_event.loc[:, "Zvuk_Reportovan_Legenda"] = df_event["Zvuk_Reportovan_Bool"].map(
        {True: "Ano", False: "Ne"}
    )
else:
    df_event.loc[:, "Zvuk_Reportovan_Text_Full"] = "Nezadáno (info chybí)"
    df_event.loc[:, "Zvuk_Reportovan_Legenda"] = "Nezadáno"
print("Předzpracování kategorií dokončeno.")


def sort_ems_key(ems_string):
    roman_map = {
        "I": 1,
        "II": 2,
        "III": 3,
        "IV": 4,
        "V": 5,
        "VI": 6,
        "VII": 7,
        "VIII": 8,
        "IX": 9,
        "X": 10,
    }
    parts = str(ems_string).split(" - ")[0]
    value = 900
    if parts in roman_map:
        value = roman_map[parts]
    elif parts == "Neklasifikováno":
        value = 1000
    return value


# --- UPRAVENÁ Funkce pro odhad EMS-98 intenzity ---
def assign_ems_intensity(row):
    popis_pohybu = str(row.get(COL_TREMOR_TYPE, "")).strip().lower()
    felt_by_str = str(row.get(COL_FELT_BY, "")).strip().lower()
    in_building_str = str(row.get(COL_IN_BUILDING, "")).strip().lower()
    fear_val = pd.to_numeric(row.get(COL_FEAR), errors="coerce")
    damage_overall_str = str(row.get(COL_DAMAGE_OVERALL, "")).strip().lower()
    in_building = in_building_str == "budova"
    felt_by_only_respondent = felt_by_str == "pouze vy"
    felt_by_some_or_many = felt_by_str in ["několik", "většina ano"]
    felt_by_most = felt_by_str == "většina ano"
    fear = fear_val == 1
    damage_overall = damage_overall_str == "bylo"

    def was_object_effect_observed(col_name_list):
        if not isinstance(col_name_list, list):
            col_name_list = [col_name_list]
        for col_name in col_name_list:
            val = row.get(col_name)
            if (
                pd.notna(val)
                and str(val).strip().lower() not in NEGATIVE_OR_EMPTY_VALUES
            ):
                return True
        return False

    if damage_overall:
        return "VI - Mírně ničivé"
    if (
        fear
        and felt_by_most
        and popis_pohybu == "silné otřesy"
        and (was_object_effect_observed(["malepredmety", "nadobi"]))
    ):
        return "VI - Mírně ničivé"
    is_strong_tremor = popis_pohybu == "silné otřesy"
    significant_object_movement_V = was_object_effect_observed(
        ["malepredmety", "nadobi", "zavespredmety", "dvere", "okna"]
    )
    if (is_strong_tremor or (felt_by_most and fear)) and significant_object_movement_V:
        return "V - Silné"
    if (
        felt_by_most
        and is_strong_tremor
        and was_object_effect_observed(["malepredmety"])
    ):
        return "V - Silné"
    object_movement_IV = was_object_effect_observed(
        ["okna", "dvere", "nadobi", "zavespredmety"]
    )
    if is_strong_tremor and not fear and object_movement_IV:
        return "IV - Značně pozorované"
    if felt_by_some_or_many and in_building and object_movement_IV:
        return "IV - Značně pozorované"
    if (
        felt_by_some_or_many
        and in_building
        and popis_pohybu in ["slabé zachvění", "chvění", "houpání"]
        and was_object_effect_observed(["okna", "dvere"])
    ):
        return "IV - Značně pozorované"
    object_movement_III_hanging = was_object_effect_observed(["zavespredmety"])
    is_weak_tremor = popis_pohybu in [
        "slabé zachvění",
        "lehké chvění",
        "houpání",
        "chvění",
    ]
    if (
        felt_by_some_or_many
        and in_building
        and is_weak_tremor
        and not was_object_effect_observed(["okna", "dvere", "nadobi", "malepredmety"])
    ):
        return "III - Slabé"
    if (
        felt_by_only_respondent
        and in_building
        and is_weak_tremor
        and not fear
        and (
            not was_object_effect_observed(COLS_OBJECT_MOVEMENT_DETAILS)
            or object_movement_III_hanging
        )
    ):
        return "III - Slabé"
    if (
        object_movement_III_hanging
        and not is_strong_tremor
        and not fear
        and not was_object_effect_observed(
            ["okna", "dvere", "nadobi", "malepredmety", "nabytektezky"]
        )
    ):
        return "III - Slabé"
    if (
        felt_by_only_respondent
        and in_building
        and popis_pohybu == "slabé zachvění"
        and not fear
        and not was_object_effect_observed(COLS_OBJECT_MOVEMENT_DETAILS)
    ):
        return "II - Zřídka pocítěno"
    if popis_pohybu == "žádný":
        if not (felt_by_some_or_many or felt_by_only_respondent):
            return "I - Nepocítěno"
        if (
            felt_by_only_respondent
            and not was_object_effect_observed(COLS_OBJECT_MOVEMENT_DETAILS)
            and not fear
        ):
            return "I - Nepocítěno"
    if popis_pohybu == "nepocítěno":
        return "I - Nepocítěno"
    return "Neklasifikováno"


df_event["EMS_Intensity_Est"] = df_event.apply(assign_ems_intensity, axis=1)
print("\n--- Odhadovaná EMS-98 Intenzita (po revizi) ---")
ems_counts = df_event["EMS_Intensity_Est"].value_counts()
sorted_ems_keys = sorted(ems_counts.index, key=lambda x: sort_ems_key(x))
ems_counts_sorted = ems_counts.reindex(sorted_ems_keys)
print(ems_counts_sorted)
if not df_event.empty:
    ems_percentages = (ems_counts_sorted / len(df_event)) * 100
    print("\nProcentuálně:")
    print(ems_percentages.round(1).astype(str) + "%")

ems_color_map = {
    "I - Nepocítěno": "rgb(200,220,255)",
    "II - Zřídka pocítěno": "rgb(160,200,255)",
    "III - Slabé": "rgb(100,220,220)",
    "IV - Značně pozorované": "rgb(120,255,120)",
    "V - Silné": "rgb(255,255,100)",
    "VI - Mírně ničivé": "rgb(255,180,100)",
    "Neklasifikováno": "rgb(200,200,200)",
}

# --- Hlavní mapa pozorování ---
print("\n--- Hlavní mapa pozorování ---")
hover_data_main_map_cols = [
    COL_OBS_DATETIME,
    "Mist_Pozorovani_Kat_Full",
    "Pocit_Kategorie_Text_Full",
    "Intenzita_Kat_Full",
    "Strach_Pocit_Kat_Text_Full",
    "Pohyb_Predmetu_Agregovany_Text_Full",
    "Poskozeni_Obecne_Text_Full",
    "Zvuk_Reportovan_Text_Full",
    "EMS_Intensity_Est",
]
valid_hover_cols = [col for col in hover_data_main_map_cols if col in df_event.columns]
main_map_png_path = create_custom_map(
    df_event,
    None,
    "Přehled pozorování",
    "pozorovani_hlavni",
    hover_data_extra=valid_hover_cols,
)
generated_map_files_for_pptx = []
if main_map_png_path:
    generated_map_files_for_pptx.append((main_map_png_path, "Přehled pozorování"))

# --- TEXTOVÉ ANALÝZY A GRAFY ---
print(f"\n--- Pozorování doma vs. venku ---")
misto_counts = df_event["Mist_Pozorovani_Kat_Full"].value_counts()
print(misto_counts)
if not df_event.empty:
    misto_percentages = (misto_counts / len(df_event)) * 100
    print("\nProcentuálně:")
    print(misto_percentages.round(1).astype(str) + "%")
create_bar_chart(
    misto_counts,
    "Pozorování doma vs. venku",
    "pozorovani_misto",
    "Místo pozorování",
    show_percentages=True,
)  # ZDE JE PRVNÍ VOLÁNÍ

# ... (zbytek kódu pro další grafy a mapy) ...
print(f"\n--- Typ pocítění ---")
pocit_counts = df_event["Pocit_Kategorie_Text_Full"].value_counts()
print(pocit_counts)
if not df_event.empty:
    pocit_percentages = (pocit_counts / len(df_event)) * 100
    print("\nProcentuálně:")
    print(pocit_percentages.round(1).astype(str) + "%")
create_bar_chart(
    pocit_counts,
    "Typ pocítění",
    "pocit_kdo",
    "Kategorie pocítění",
    show_percentages=True,
)
print(f"\n--- Intenzita otřesů (popis) ---")
intenzita_counts = df_event["Intenzita_Kat_Full"].value_counts()
print(intenzita_counts)
if not df_event.empty:
    intenzita_percentages = (intenzita_counts / len(df_event)) * 100
    print("\nProcentuálně:")
    print(intenzita_percentages.round(1).astype(str) + "%")
create_bar_chart(
    intenzita_counts,
    "Intenzita otřesů (popis)",
    "intenzita_popis",
    "Popis intenzity",
    show_percentages=True,
)
print(f"\n--- Strach/Panika ---")
strach_counts = df_event["Strach_Pocit_Kat_Text_Full"].value_counts()
print(strach_counts)
if not df_event.empty:
    strach_percentages = (strach_counts / len(df_event)) * 100
    print("\nProcentuálně:")
    print(strach_percentages.round(1).astype(str) + "%")
create_bar_chart(
    strach_counts,
    "Pocit strachu/paniky",
    "strach_panika",
    "Hlášení strachu/paniky",
    show_percentages=True,
)
print("\n--- Pohyb předmětů ---")
if "Pohyb_Predmetu_Agregovany_Text_Full" in df_event.columns:
    pohyb_agreg_counts = df_event["Pohyb_Predmetu_Agregovany_Text_Full"].value_counts()
    print(pohyb_agreg_counts)
    if not df_event.empty:
        pohyb_agreg_percentages = (pohyb_agreg_counts / len(df_event)) * 100
        print("\nProcentuálně:")
        print(pohyb_agreg_percentages.round(1).astype(str) + "%")
    if pohyb_agreg_counts.get("Ano (pohyb předmětů)", 0) > 0:
        create_bar_chart(
            pohyb_agreg_counts,
            "Agregovaný pohyb předmětů",
            "pohyb_predmetu_agreg",
            "Pozorován pohyb?",
            show_percentages=True,
        )
    if actual_movement_detail_cols:
        movement_details_data = {}
        for col_name in actual_movement_detail_cols:
            s = df_event[col_name]
            s_str_lower_stripped = s.astype(str).str.strip().str.lower()
            condition = s.notna() & (
                ~s_str_lower_stripped.isin(NEGATIVE_OR_EMPTY_VALUES)
            )
            observed_count = condition.sum()
            if observed_count > 0:
                movement_details_data[col_name] = observed_count
        if movement_details_data:
            pohyb_detail_series = pd.Series(movement_details_data).sort_values(
                ascending=False
            )
            print("Detaily pohybů (počet):")
            print(pohyb_detail_series)
            if not df_event.empty:
                pohyb_detail_percentages = (pohyb_detail_series / len(df_event)) * 100
                print("\nProcentuálně (z celkového počtu pozorování):")
                print(pohyb_detail_percentages.round(1).astype(str) + "%")
            create_bar_chart(
                pohyb_detail_series,
                "Detaily pohybů předmětů",
                "pohyb_detaily",
                "Typ pohybu",
                show_percentages=True,
            )
print(f"\n--- Poškození budov ---")
poskozeni_counts = df_event["Poskozeni_Obecne_Text_Full"].value_counts()
print(poskozeni_counts)
if not df_event.empty:
    poskozeni_percentages = (poskozeni_counts / len(df_event)) * 100
    print("\nProcentuálně:")
    print(poskozeni_percentages.round(1).astype(str) + "%")
create_bar_chart(
    poskozeni_counts,
    "Poškození budov",
    "poskozeni_budov",
    "Poškození hlášeno?",
    show_percentages=True,
)
print(f"\n--- Analýza Zvuků ---")
if "Zvuk_Reportovan_Text_Full" in df_event.columns:
    zvuk_agreg_counts = df_event["Zvuk_Reportovan_Text_Full"].value_counts()
    print(zvuk_agreg_counts)
    if not df_event.empty:
        zvuk_agreg_percentages = (zvuk_agreg_counts / len(df_event)) * 100
        print("\nProcentuálně:")
        print(zvuk_agreg_percentages.round(1).astype(str) + "%")
    if zvuk_agreg_counts.get("Ano (zvuk reportován)", 0) > 0:
        create_bar_chart(
            zvuk_agreg_counts,
            "Agregovaný report zvuků",
            "zvuky_agreg",
            "Zvuk reportován?",
            show_percentages=True,
        )
    if actual_sound_cols:
        sound_details_data = {}
        for col_name in actual_sound_cols:
            s = df_event[col_name]
            s_str_lower_stripped = s.astype(str).str.strip().str.lower()
            condition = s.notna() & (
                ~s_str_lower_stripped.isin(NEGATIVE_OR_EMPTY_VALUES)
            )
            observed_count = condition.sum()
            if observed_count > 0:
                sound_details_data[col_name] = observed_count
        if sound_details_data:
            zvuk_detail_series = pd.Series(sound_details_data).sort_values(
                ascending=False
            )
            print("Detaily zvuků (počet):")
            print(zvuk_detail_series)
            if not df_event.empty:
                zvuk_detail_percentages = (zvuk_detail_series / len(df_event)) * 100
                print("\nProcentuálně (z celkového počtu pozorování):")
                print(zvuk_detail_percentages.round(1).astype(str) + "%")
            create_bar_chart(
                zvuk_detail_series,
                "Detaily reportovaných zvuků",
                "zvuky_detaily",
                "Typ zvuku",
                show_percentages=True,
            )

print("\n--- Generování parametrických map ---")
map_configs = [
    (
        "Pocit_Kategorie_Legenda",
        "Kdo pocítil otřesy",
        "felt_by",
        {"Většina": "green", "Jen respondent": "orange", "Nezadáno": "lightgrey"},
        None,
        "Pocit_Kategorie_Text_Full",
    ),
    (
        "Zvuk_Reportovan_Legenda",
        "Reportované zvuky",
        "sounds_reported",
        {"Ano": "purple", "Ne": "lightskyblue", "Nezadáno": "lightgrey"},
        None,
        "Zvuk_Reportovan_Text_Full",
    ),
    (
        "Poskozeni_Obecne_Legenda",
        "Poškození budov",
        "damage",
        {"Ano": "darkred", "Ne": "darkgreen", "Nezadáno": "lightgrey"},
        None,
        "Poskozeni_Obecne_Text_Full",
    ),
    (
        "Pohyb_Predmetu_Legenda",
        "Pohyb předmětů",
        "object_movement",
        {"Ano": "orangered", "Ne": "mediumseagreen", "Nezadáno": "lightgrey"},
        None,
        "Pohyb_Predmetu_Agregovany_Text_Full",
    ),
    (
        "Intenzita_Kat_Legenda",
        "Intenzita otřesů (popis)",
        "tremor_intensity_descr",
        None,
        None,
        "Intenzita_Kat_Full",
    ),
    (
        "Strach_Pocit_Legenda",
        "Pocit strachu/paniky",
        "fear",
        {"Ano": "crimson", "Ne/Nezadáno": "teal"},
        None,
        "Strach_Pocit_Kat_Text_Full",
    ),
    (
        "Mist_Pozorovani_Legenda",
        "Místo pozorování",
        "in_building",
        {"Doma": "sandybrown", "Venku/Nezadáno": "skyblue"},
        None,
        "Mist_Pozorovani_Kat_Full",
    ),
]
for config_idx, config in enumerate(map_configs):
    col_for_color, title, fname, cmap, hover_extra, col_for_hover = config
    cat_order_current = []
    if col_for_color in df_event.columns and df_event[col_for_color].nunique() > 0:
        unique_values = df_event[col_for_color].unique().tolist()
        if cmap:
            cat_order_current = list(cmap.keys())
            missing = sorted([v for v in unique_values if v not in cat_order_current])
            cat_order_current.extend(missing)
        else:
            nezadano_like = [
                "Nezadáno",
                "Nezadáno/Jiné",
                "Nezadáno (info chybí)",
                "Neklasifikováno",
                "nan",
            ]
            standard_vals = sorted(
                [
                    v
                    for v in unique_values
                    if str(v) not in nezadano_like and not pd.isna(v)
                ],
                key=lambda x: str(x).lower(),
            )
            nezadano_vals = sorted(
                [v for v in unique_values if str(v) in nezadano_like or pd.isna(v)],
                key=lambda x: str(x).lower(),
            )
            cat_order_current = standard_vals + nezadano_vals
    png_path = create_custom_map(
        df_event,
        col_for_color,
        title,
        fname,
        {col_for_color: cat_order_current}
        if cat_order_current and col_for_color in df_event.columns
        else None,
        cmap,
        hover_extra,
        col_for_hover,
        show_isoseismal_areas=False,
    )
    if png_path:
        generated_map_files_for_pptx.append((png_path, title))

print("\n--- Generování EMS mapy s izoseismálními oblastmi ---")
ems_cat_order = sorted(df_event["EMS_Intensity_Est"].unique(), key=sort_ems_key)
ems_hulls_map_title = "Odhad EMS-98 Intenzita s oblastmi"
ems_hulls_map_png_path = create_custom_map(
    df_event,
    "EMS_Intensity_Est",
    ems_hulls_map_title,
    "ems_intensity_hulls",
    {"EMS_Intensity_Est": ems_cat_order},
    ems_color_map,
    [
        COL_TREMOR_TYPE,
        COL_FELT_BY,
        COL_DAMAGE_OVERALL,
        "Pohyb_Predmetu_Agregovany_Text_Full",
    ],
    "EMS_Intensity_Est",
    show_isoseismal_areas=True,
    ems_color_map_for_hulls=ems_color_map,
    sort_ems_key_func=sort_ems_key,
)
if ems_hulls_map_png_path:
    generated_map_files_for_pptx.append((ems_hulls_map_png_path, ems_hulls_map_title))

print("\n--- Generování PowerPoint prezentace ---")
if generated_map_files_for_pptx:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    blank_slide_layout = prs.slide_layouts[6]
    for png_path, map_slide_title in generated_map_files_for_pptx:
        if os.path.exists(png_path):
            slide = prs.slides.add_slide(blank_slide_layout)
            title_shape = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)
            )
            title_frame = title_shape.text_frame
            title_frame.text = f"{EQ_LOCATION_NAME}: {map_slide_title}"
            title_frame.paragraphs[0].font.size = Inches(0.24)
            title_frame.paragraphs[0].font.bold = True
            img_width_on_slide = Inches(9)
            img_height_on_slide = img_width_on_slide * (750 / 1000)
            left = (prs.slide_width - img_width_on_slide) / 2
            top = Inches(0.75)
            try:
                slide.shapes.add_picture(
                    png_path,
                    left,
                    top,
                    width=img_width_on_slide,
                    height=img_height_on_slide,
                )
                print(f"Přidána mapa '{map_slide_title}' do prezentace.")
            except Exception as e:
                print(f"CHYBA při přidávání obrázku {png_path} do prezentace: {e}")
        else:
            print(f"VAROVÁNÍ: Soubor s mapou {png_path} nebyl nalezen.")
    pptx_filename = os.path.join(
        OUTPUT_DIR,
        f"prezentace_mapy_{EQ_LOCATION_NAME.lower().replace(' ', '_')}_{EQ_YEAR_TARGET}.pptx",
    )
    try:
        prs.save(pptx_filename)
        print(f"PowerPoint prezentace uložena do: {pptx_filename}")
    except Exception as e:
        print(f"CHYBA při ukládání PowerPoint prezentace: {e}")
else:
    print("Nebyly vygenerovány žádné mapy pro přidání do PowerPoint prezentace.")

print("\n--- SKRIPT ÚSPĚŠNĚ DOKONČEN ---")
